from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import date

COMPANY   = "Janatics India Private Limited"
CIN       = "U31103TZ1991PTC003409"
DATE_STR  = date.today().strftime("%d %B %Y")
COMPOSITE = 64

DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x00, 0xB0, 0x50)
AMBER      = RGBColor(0xFF, 0xC0, 0x00)
RED_C      = RGBColor(0xFF, 0x00, 0x00)
BLUE_FR    = RGBColor(0x25, 0x63, 0xEB)
DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)

def maturity_label(s):
    if s >= 84: return "Future Ready"
    if s >= 67: return "Strategic"
    if s >= 34: return "Siloed"
    return "Legacy"

def score_color(s):
    if s >= 84: return BLUE_FR
    if s >= 67: return GREEN
    if s >= 34: return AMBER
    return RED_C

def add_text_box(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=None,
                 bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]

# ── Slide 1: Title ─────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, DARK_BLUE)
add_rect(slide, 0, Inches(4.8), W, Inches(2.7), MID_BLUE)

add_text_box(slide, "DIGITAL DIAGNOSTIC AND MATURITY REPORT",
             Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.8),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, COMPANY,
             Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.9),
             font_size=22, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text_box(slide, f"CIN: {CIN}",
             Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.5),
             font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.9), Inches(3.8), Inches(3.5), Inches(1.2), score_color(COMPOSITE))
add_text_box(slide, f"{COMPOSITE}/100",
             Inches(4.9), Inches(3.8), Inches(3.5), Inches(0.7),
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, f"COMPOSITE DDMR SCORE — {maturity_label(COMPOSITE)}",
             Inches(4.9), Inches(4.5), Inches(3.5), Inches(0.4),
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, f"Report Date: {DATE_STR}   |   Confidential",
             Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
             font_size=10, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# ── Slide 2: Company Snapshot ──────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text_box(slide, "Company Snapshot",
             Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8),
             font_size=20, bold=True, color=WHITE)

facts = [
    ("Founded",           "1977 (brand) / 1991 (incorporated)"),
    ("CIN",               "U31103TZ1991PTC003409"),
    ("Headquarters",      "Coimbatore, Tamil Nadu"),
    ("Additional Plants", "Ahmedabad, Noida"),
    ("Status",            "Active | ROC Coimbatore"),
    ("Revenue (FY2025)",  "₹531 Crore"),
    ("Employees",         "662"),
    ("Products",          "3,500+ SKUs"),
    ("Customers",         "17,000+"),
    ("Countries",         "42+"),
    ("Subsidiaries",      "USA, Germany"),
    ("Credit Rating",     "ICRA — Stable Outlook"),
    ("R&D",               "DSIR-recognized"),
    ("Industry",          "Pneumatics & Industrial Automation"),
]

cols = [facts[:7], facts[7:]]
for ci, col in enumerate(cols):
    for ri, (label, val) in enumerate(col):
        lx = Inches(0.4 + ci * 6.5)
        ty = Inches(1.3 + ri * 0.82)
        add_rect(slide, lx, ty, Inches(2.3), Inches(0.55), MID_BLUE)
        add_text_box(slide, label, lx + Inches(0.05), ty + Inches(0.05),
                     Inches(2.2), Inches(0.45), font_size=9, bold=True, color=WHITE)
        add_rect(slide, lx + Inches(2.35), ty, Inches(3.8), Inches(0.55), LIGHT_BLUE)
        add_text_box(slide, val, lx + Inches(2.4), ty + Inches(0.05),
                     Inches(3.7), Inches(0.45), font_size=10, color=DARK_BLUE)

# ── Slide 3: Scorecard Radar ───────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text_box(slide, "DDMR Digital Maturity Scorecard",
             Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8),
             font_size=20, bold=True, color=WHITE)

chart_data = ChartData()
chart_data.categories = [
    "Strategy",
    "Operations &\nSupply Chain",
    "Sales &\nMarketing",
    "Technology\nAdoption",
    "Skills &\nCapabilities",
]
chart_data.add_series("Score", (70, 65, 55, 70, 60))
chart_data.add_series("Benchmark (67)", (67, 67, 67, 67, 67))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.RADAR_FILLED,
    Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.8),
    chart_data
).chart
chart.has_legend = True

dims_summary = [
    ("ST",  "Strategy",                  70),
    ("OSC", "Operations & Supply Chain", 65),
    ("SM",  "Sales & Marketing",         55),
    ("TA",  "Technology Adoption",       70),
    ("SC",  "Skills & Capabilities",     60),
]
for i, (code, name, score) in enumerate(dims_summary):
    ty = Inches(1.4 + i * 1.1)
    add_rect(slide, Inches(8.3), ty, Inches(0.6), Inches(0.7), score_color(score))
    add_text_box(slide, str(score), Inches(8.3), ty, Inches(0.6), Inches(0.7),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(8.95), ty, Inches(4.0), Inches(0.7), LIGHT_BLUE)
    add_text_box(slide, f"{code}: {name}  [{maturity_label(score)}]",
                 Inches(9.0), ty + Inches(0.08),
                 Inches(3.9), Inches(0.55), font_size=10, color=DARK_BLUE)

add_rect(slide, Inches(8.3), Inches(7.0), Inches(4.65), Inches(0.45), score_color(COMPOSITE))
add_text_box(slide, f"Composite Score: {COMPOSITE}/100  —  {maturity_label(COMPOSITE)}",
             Inches(8.3), Inches(7.0), Inches(4.65), Inches(0.45),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slides 4–8: One per dimension ─────────────────────────────────────────────

dimension_slides = [
    {
        "code": "ST", "name": "Strategy", "score": 70, "weight": "20%",
        "bullets": [
            "Product portfolio pivot to I4.0 equipment — electric actuators, smart sensors, robotics",
            "DSIR-recognized R&D division — structured technology investment commitment",
            "No formal digital transformation roadmap or DX strategy publicly disclosed",
            "No chief digital officer or digital leadership appointment evidenced",
            "Technology partnerships limited to DSIR; no digital ecosystem alliances disclosed",
            "Innovation pipeline: incremental product extension; no open innovation signals",
        ],
        "sub_scores": [("ST01 DX Strategy",65),("ST02 Leadership Vision",72),("ST03 Investment",75),
                       ("ST04 Tech Partnerships",68),("ST05 Innovation",70)],
        "highlight": "DSIR R&D + I4.0 product pivot signal digital leadership awareness",
        "risk": "No formal DX roadmap or strategy publicly articulated",
    },
    {
        "code": "OSC", "name": "Operations & Supply Chain", "score": 65, "weight": "20%",
        "bullets": [
            "70,000 sqm Coimbatore HQ | Ahmedabad + Noida plants — strong physical footprint",
            "CNC machining, casting, surface treatment — physical automation well established",
            "No ERP, MES, or supply chain visibility platform publicly disclosed",
            "No Industry 4.0 certification or smart factory initiative announced",
            "I4.0 product portfolio demonstrates concept awareness but internal adoption unclear",
            "Quality management practices consistent with ICRA compliance; formal system unconfirmed",
        ],
        "sub_scores": [("OSC01 Mfg Readiness",60),("OSC02 SC Visibility",50),("OSC03 Automation",70),
                       ("OSC04 Quality Mgmt",72),("OSC05 I4.0 Integration",73)],
        "highlight": "Strong physical automation; CNC, casting, 3 plant locations",
        "risk": "No ERP/MES/SCM platform or I4.0 certification evidenced publicly",
    },
    {
        "code": "SM", "name": "Sales & Marketing", "score": 55, "weight": "20%",
        "bullets": [
            "Two product websites (janatics.com / janaticspneumatics.com) — functional, not digital-first",
            "17,000+ customers across 7 industries — managed at scale; CRM tool undisclosed",
            "No B2B e-commerce, digital ordering portal, or online pricing identified",
            "Limited LinkedIn / social media engagement signals from public sources",
            "No paid digital marketing, content strategy, or SEO investment evidenced",
            "200 global distribution partners — relationship-driven, not channel-digital",
        ],
        "sub_scores": [("SM01 Digital Presence",60),("SM02 Social Media",40),("SM03 E-commerce",45),
                       ("SM04 Digital Marketing",55),("SM05 CRM / Data",75)],
        "highlight": "17,000+ customers managed; CRM implied at scale",
        "risk": "No digital sales channel, e-commerce, or social marketing strategy",
    },
    {
        "code": "TA", "name": "Technology Adoption", "score": 70, "weight": "20%",
        "bullets": [
            "DSIR R&D developing robotic systems, IoT sensors, electric actuators — strong product R&D",
            "Industry 4.0 didactic product line — smart manufacturing expertise in products",
            "Website infrastructure functional; no modern CMS, API layer, or customer portal",
            "No cloud migration, SaaS adoption, or digital collaboration tools disclosed",
            "Cybersecurity: ICRA-compliant; no ISO 27001 or public data policy found",
            "Largest gap: internal tech infrastructure lags product-layer digital maturity",
        ],
        "sub_scores": [("TA01 Digital Infra",65),("TA02 IoT Products",80),("TA03 Cloud/SaaS",60),
                       ("TA04 R&D Digital",78),("TA05 Cybersecurity",67)],
        "highlight": "IoT product R&D (80/100) — DSIR-backed, robotic and sensor development",
        "risk": "Internal cloud/SaaS adoption and cybersecurity framework not evidenced",
    },
    {
        "code": "SC", "name": "Skills & Capabilities", "score": 60, "weight": "20%",
        "bullets": [
            "662 employees — engineering-deep workforce; strong CNC and automation specialists",
            "DSIR R&D team in place — structured technical innovation capability",
            "Limited signals of software engineers, data scientists, or digital roles hired",
            "No public L&D, digital upskilling, or online learning program disclosed",
            "Management I4.0 awareness evidenced by product decisions; no CDO hire signal",
            "No hackathon, innovation lab, startup engagement, or employee innovation program",
        ],
        "sub_scores": [("SC01 Digital Talent",55),("SC02 Tech Workforce",72),("SC03 Training",50),
                       ("SC04 Leadership DQ",65),("SC05 Innovation Culture",58)],
        "highlight": "Engineering-depth workforce; DSIR R&D team in place",
        "risk": "No digital talent pipeline or upskilling program evidenced publicly",
    },
]

for d in dimension_slides:
    slide = prs.slides.add_slide(blank_layout)

    add_rect(slide, 0, 0, W, Inches(1.1), score_color(d["score"]))
    add_text_box(slide, f"{d['code']}: {d['name']}",
                 Inches(0.3), Inches(0.1), Inches(9.0), Inches(0.55),
                 font_size=20, bold=True, color=WHITE)
    add_text_box(slide, f"Score: {d['score']}/100   |   Maturity: {maturity_label(d['score'])}   |   Weight: {d['weight']}",
                 Inches(0.3), Inches(0.6), Inches(9.0), Inches(0.4),
                 font_size=12, color=WHITE)

    add_rect(slide, Inches(10.8), Inches(0.1), Inches(2.2), Inches(0.9), DARK_BLUE)
    add_text_box(slide, f"{d['score']}/100",
                 Inches(10.8), Inches(0.1), Inches(2.2), Inches(0.9),
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for bi, bullet in enumerate(d["bullets"]):
        by = Inches(1.25 + bi * 0.87)
        add_rect(slide, Inches(0.25), by + Inches(0.15), Inches(0.08), Inches(0.22),
                 score_color(d["score"]))
        add_text_box(slide, bullet,
                     Inches(0.45), by, Inches(8.1), Inches(0.75),
                     font_size=11, color=DARK_GRAY)

    panel_x = Inches(8.8)
    add_rect(slide, panel_x, Inches(1.1), Inches(4.3), Inches(0.4), MID_BLUE)
    add_text_box(slide, "Sub-Metric Scores",
                 panel_x, Inches(1.1), Inches(4.3), Inches(0.4),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bar_h = Inches(0.52)
    for si, (sm_name, sm_score) in enumerate(d["sub_scores"]):
        sy = Inches(1.55 + si * 0.65)
        add_text_box(slide, sm_name, panel_x, sy, Inches(2.0), bar_h,
                     font_size=9, color=DARK_BLUE)
        add_rect(slide, panel_x + Inches(2.05), sy + Inches(0.12),
                 Inches(1.8), Inches(0.28), LIGHT_BLUE)
        bar_w = Inches(1.8 * sm_score / 100)
        add_rect(slide, panel_x + Inches(2.05), sy + Inches(0.12),
                 bar_w, Inches(0.28), score_color(sm_score))
        add_text_box(slide, str(sm_score),
                     panel_x + Inches(3.9), sy, Inches(0.4), bar_h,
                     font_size=9, bold=True, color=score_color(sm_score))

    add_rect(slide, 0, Inches(6.85), W * 0.6, Inches(0.55), GREEN)
    add_text_box(slide, f"Key Strength: {d['highlight']}",
                 Inches(0.15), Inches(6.87), Inches(7.8), Inches(0.45),
                 font_size=9, bold=True, color=WHITE)
    add_rect(slide, W * 0.6, Inches(6.85), W * 0.4, Inches(0.55), RED_C)
    add_text_box(slide, f"Risk: {d['risk']}",
                 W * 0.6 + Inches(0.1), Inches(6.87), Inches(5.1), Inches(0.45),
                 font_size=9, bold=True, color=WHITE)

# ── Slide 9: Next Best Actions ────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text_box(slide, "Next Best Actions — Digital Transformation Roadmap",
             Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8),
             font_size=20, bold=True, color=WHITE)

next_actions_deck = [
    {
        "rank": "#01  HIGH",
        "color": RED_C,
        "title": "Activate B2B Digital Sales Channel",
        "gap": "SM03: 45/100",
        "benefit": "+15–20% incremental revenue from direct digital orders within 18 months",
        "loss": "Global competitors (SMC, Festo) own digital search — every quarter cedes buyer mindshare",
        "why": "India B2B e-commerce growing 25%+ p.a.; 49-yr brand + 3,500 SKUs = digital-first advantage",
    },
    {
        "rank": "#02  HIGH",
        "color": RED_C,
        "title": "Deploy ERP + Supply Chain Visibility",
        "gap": "OSC02: 50/100",
        "benefit": "10–15% operational efficiency gain; unlock global OEM export qualifications",
        "loss": "OEM buyers mandate digital traceability — no ERP = disqualified from high-value RFQs",
        "why": "China+1 OEM inquiry window is open NOW — digital ops credibility needed to convert",
    },
    {
        "rank": "#03  MEDIUM",
        "color": AMBER,
        "title": "Digital Workforce Upskilling Program",
        "gap": "SC03: 50/100 · SC01: 55/100",
        "benefit": "Accelerates all digital tool adoption; reduces transformation failure risk by 60%+",
        "loss": "Technology without people = stranded investment; ERP + e-commerce will underperform",
        "why": "Digital skills gap widening; early movers gain 2–3 yr head start on talent",
    },
    {
        "rank": "#04  MEDIUM",
        "color": AMBER,
        "title": "Smart Factory Showcase — Coimbatore Plant",
        "gap": "TA01: 65/100 · ST03: 75/100",
        "benefit": "Credibility proof-by-example; accelerates I4.0 sales cycles; OEM audit readiness",
        "loss": "Selling I4.0 while running a non-digitized plant is a credibility gap buyers detect",
        "why": "Janatics owns the tech — fastest credibility play with no new vendor or capex risk",
    },
]

col_w = Inches(3.1)
gap = Inches(0.13)
start_x = Inches(0.3)

for ci, item in enumerate(next_actions_deck):
    x = start_x + ci * (col_w + gap)
    add_rect(slide, x, Inches(1.15), col_w, Inches(0.42), item["color"])
    add_text_box(slide, item["rank"], x, Inches(1.15), col_w, Inches(0.42),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(1.60), col_w, Inches(0.52), MID_BLUE)
    add_text_box(slide, item["title"], x + Inches(0.05), Inches(1.62),
                 col_w - Inches(0.1), Inches(0.48),
                 font_size=11, bold=True, color=WHITE)
    add_rect(slide, x, Inches(2.15), col_w, Inches(0.32), LIGHT_BLUE)
    add_text_box(slide, f"Gap: {item['gap']}", x + Inches(0.05), Inches(2.16),
                 col_w - Inches(0.1), Inches(0.30),
                 font_size=8, bold=True, color=DARK_BLUE)

    row_items = [
        ("Expected Benefit", item["benefit"], RGBColor(0xD5, 0xE8, 0xD4), RGBColor(0x1F, 0x38, 0x64)),
        ("Cost of Inaction", item["loss"],    RGBColor(0xF8, 0xCE, 0xCC), RGBColor(0x7F, 0x00, 0x00)),
        ("Why Now",          item["why"],     RGBColor(0xDA, 0xE8, 0xFC), RGBColor(0x1F, 0x38, 0x64)),
    ]
    for ri, (label, text, bg, tc) in enumerate(row_items):
        ry = Inches(2.50 + ri * 1.55)
        add_rect(slide, x, ry, col_w, Inches(0.30), item["color"] if ri == 0 else (
            RED_C if ri == 1 else MID_BLUE))
        add_text_box(slide, label, x + Inches(0.05), ry, col_w, Inches(0.30),
                     font_size=8, bold=True, color=WHITE)
        add_rect(slide, x, ry + Inches(0.30), col_w, Inches(1.22), bg)
        add_text_box(slide, text, x + Inches(0.06), ry + Inches(0.33),
                     col_w - Inches(0.12), Inches(1.18),
                     font_size=8, color=tc)

add_text_box(slide, "Maturity band target after initiatives: Siloed → Strategic (67–83)",
             Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             bg_color=MID_BLUE)

# ── Slide 10: Recommendation ───────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, Inches(1.1), DARK_BLUE)
add_text_box(slide, "Executive Recommendation",
             Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8),
             font_size=20, bold=True, color=WHITE)

add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.7), score_color(COMPOSITE))
add_text_box(slide, f"DDMR Composite Score: {COMPOSITE}/100 — {maturity_label(COMPOSITE).upper()}  |  Digital Transformation: Early-Stage Engagement Opportunity",
             Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.7),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rec_points = [
    ("Strengths to leverage",
     "DSIR R&D and I4.0 product portfolio demonstrate product-layer digital competence. "
     "Strong manufacturing base (70k sqm, 3 plants, 3,500+ SKUs) and ICRA Stable credit "
     "provide a stable foundation for digital investment."),
    ("Priority gaps",
     "Digital sales channel (SM03: 45/100) and supply chain visibility (OSC02: 50/100) are "
     "the two lowest-scoring sub-metrics — quick wins with measurable ROI. Social media "
     "engagement (SM02: 40/100) requires immediate activation."),
    ("Recommended DX initiatives",
     "1. Deploy B2B digital ordering portal and SEO-optimised product catalog. "
     "2. Publish ERP/SCM roadmap — signals operational maturity to global OEM buyers. "
     "3. Launch digital upskilling program for engineering workforce (SC03: 50/100). "
     "4. Appoint digital transformation lead or CDO."),
    ("Opportunity",
     "Janatics' own I4.0 product expertise creates a credibility-led path to internal "
     "digitization — using their own sensors and automation products to modernize their "
     "manufacturing is a compelling narrative for global OEM customers assessing vendor "
     "digital maturity."),
]

for ri, (heading, body) in enumerate(rec_points):
    ty = Inches(2.1 + ri * 1.3)
    add_rect(slide, Inches(0.4), ty, Inches(2.8), Inches(1.1), MID_BLUE)
    add_text_box(slide, heading, Inches(0.45), ty + Inches(0.1),
                 Inches(2.7), Inches(0.9), font_size=10, bold=True, color=WHITE)
    add_rect(slide, Inches(3.3), ty, Inches(9.6), Inches(1.1), LIGHT_BLUE)
    add_text_box(slide, body, Inches(3.4), ty + Inches(0.05),
                 Inches(9.4), Inches(1.0), font_size=9, color=DARK_BLUE)

add_text_box(slide, f"Report generated: {DATE_STR}  |  {COMPANY}  |  DDMR v1.0",
             Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
             font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER, italic=True)

out = "/Users/laaksandy/Documents/Claude-testing/DDMR/output/Janatics_DDMR_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
